"""Build a GraphRAG-Bench-style dataset from local DOCX/PPTX brand reference docs."""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import zipfile
import xml.etree.ElementTree as ET
from collections import defaultdict
from pathlib import Path

import pandas as pd
from docx import Document
from pptx import Presentation

from .. import config


QUESTION_TYPES = {
    "fact": "Fact Retrieval",
    "reasoning": "Complex Reasoning",
    "summary": "Contextual Summarize",
    "creative": "Creative Generation",
}

SPACE_RE = re.compile(r"\s+")
BULLET_RE = re.compile(r"^[\u2022\-\*\u25cf\u25aa\u25e6\u2013\u2014]+\s*")


def _clean(text: str) -> str:
    text = (text or "").replace("\x0b", "\n").replace("\xa0", " ")
    lines = []
    for line in text.splitlines():
        line = BULLET_RE.sub("", line.strip())
        line = SPACE_RE.sub(" ", line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def _one_line(text: str) -> str:
    return SPACE_RE.sub(" ", _clean(text)).strip()


def _norm(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", _one_line(text).lower()).strip()


def _title(path: Path) -> str:
    title = re.sub(r"[_]+", " ", path.stem)
    title = re.sub(r"\s*\(\d+\)\s*", " ", title)
    title = re.sub(r"\b(v\d+|M1|APRIL\s+2026|012026|2026)\b", " ", title, flags=re.I)
    return SPACE_RE.sub(" ", title).strip(" -")


def _topic(path: Path) -> str:
    """Turn filenames into natural domain topics, not document references."""
    raw = _title(path)
    low = raw.lower()
    if "a b test engagement" in low:
        return "the MCC in-stream ads engagement lift test"
    if "ad prioritisation" in low or "ad prioritization" in low:
        return "video feed ad prioritization"
    if "feed personalization" in low:
        return "feed personalization for placements"
    if "ad ops guide" in low and "mcc" in low:
        return "MCC ad operations"
    if "mastercard" in low and "technical" in low:
        return "the Mastercard technical integration"
    if "mastercard" in low and "sponsorship" in low:
        return "the Mastercard sponsorship program"
    if "bcc setup" in low:
        return "BCC setup in Octo"
    if "octo canvas" in low or "octocanvas" in low:
        return "Octo Canvas"
    if "liveramp" in low:
        return "the LiveRamp identity integration"
    if "iheart" in low and "audience extension" in low:
        return "iHeart audience extension"
    if "iheart" in low and "monetize" in low:
        return "Genuin and iHeart O&O monetization"
    if "iheart" in low and "feature proposal" in low:
        return "the iHeart and Genuin feature proposal"
    if "iheartvip" in low:
        return "the iHeartVIP revival initiative"
    if "grubhub" in low:
        return "the Grubhub Genuin sponsorship program"
    if "bettercollective" in low:
        return "the Better Collective Genuin program"
    if "publicis health media" in low or "gsk" in low:
        return "Genuin for Publicis Health Media and GSK"
    if "o&o data flow" in low:
        return "the O&O data flow"
    if "product technology raci" in low or "product & technology raci" in low:
        return "Genuin product and technology ownership"
    if "product understanding" in low:
        return "Genuin product and process improvement"
    if "user experiences" in low:
        return "Genuin user experiences"
    if "advertising on genuin for tiktok" in low:
        return "advertising on Genuin for TikTok"
    if "executive one-pager" in low or "monetize" in low:
        return "Genuin Monetize"
    if "prd" in low and "ad placement" in low:
        return "ad placement configuration"
    if "quarterly company" in low:
        return "the quarterly company update"
    if "org chart" in low:
        return "Genuin's organization structure"
    return raw


def _trim(text: str, max_chars: int = 950) -> str:
    text = _one_line(text)
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars]
    last = max(cut.rfind(". "), cut.rfind("; "), cut.rfind(", "))
    if last > 350:
        cut = cut[:last + 1]
    return cut.rstrip() + "..."


def _make_id(source: str, question: str, index: int) -> str:
    digest = hashlib.sha1(f"{source}|{question}|{index}".encode("utf-8")).hexdigest()[:10]
    return f"brand_ref_{digest}"


def _extract_docx_xml(path: Path) -> list[dict]:
    """Fallback for malformed DOCX packages that python-docx refuses to open."""
    blocks = []
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(path) as zf:
        names = [
            name for name in zf.namelist()
            if name.startswith("word/")
            and (name == "word/document.xml" or name.startswith("word/header") or name.startswith("word/footer"))
        ]
        for name in names:
            root = ET.fromstring(zf.read(name))
            for para in root.findall(".//w:p", ns):
                text = _clean("".join(t.text or "" for t in para.findall(".//w:t", ns)))
                if text:
                    blocks.append({
                        "kind": "paragraph",
                        "index": len(blocks) + 1,
                        "style": f"xml:{name}",
                        "text": text,
                    })
            for table in root.findall(".//w:tbl", ns):
                rows = []
                for row in table.findall(".//w:tr", ns):
                    vals = []
                    for cell in row.findall(".//w:tc", ns):
                        text = _one_line("".join(t.text or "" for t in cell.findall(".//w:t", ns)))
                        if text:
                            vals.append(text)
                    if vals:
                        rows.append(" | ".join(vals))
                if rows:
                    blocks.append({
                        "kind": "table",
                        "index": len(blocks) + 1,
                        "style": f"xml:{name}",
                        "text": "\n".join(rows),
                    })
    return blocks


def _extract_docx(path: Path) -> tuple[list[dict], str]:
    try:
        doc = Document(str(path))
        blocks = []
        for index, para in enumerate(doc.paragraphs, 1):
            text = _clean(para.text)
            if text:
                blocks.append({
                    "kind": "paragraph",
                    "index": index,
                    "style": getattr(para.style, "name", "") if para.style else "",
                    "text": text,
                })
        for index, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                vals = [_one_line(cell.text) for cell in row.cells]
                vals = [val for val in vals if val]
                if vals:
                    rows.append(" | ".join(vals))
            if rows:
                blocks.append({
                    "kind": "table",
                    "index": index,
                    "style": "table",
                    "text": "\n".join(rows),
                })
        return blocks, "python-docx"
    except Exception:
        return _extract_docx_xml(path), "docx-xml-fallback"


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _extract_pptx(path: Path) -> tuple[list[dict], str]:
    prs = Presentation(str(path))
    blocks = []
    for slide_number, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                text = _clean(shape.text)
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    vals = [_one_line(cell.text) for cell in row.cells]
                    vals = [val for val in vals if val]
                    if vals:
                        rows.append(" | ".join(vals))
                if rows:
                    parts.append("\n".join(rows))
        seen = set()
        unique = []
        for part in parts:
            if part not in seen:
                seen.add(part)
                unique.append(part)
        text = _clean("\n".join(unique))
        if text:
            blocks.append({
                "kind": "slide",
                "index": slide_number,
                "style": "slide",
                "text": text,
            })
    return blocks, "python-pptx"


def _is_heading(line: str) -> bool:
    if not (4 <= len(line) <= 95) or line.endswith((".", ",", ";", ":")):
        return False
    words = line.split()
    if len(words) > 12 or line.lower() in {"agenda", "thank you", "appendix", "contents"}:
        return False
    if not re.sub(r"[^A-Za-z]", "", line):
        return False
    cap_words = sum(1 for word in words if word[:1].isupper() or word.isupper())
    return cap_words >= max(1, len(words) // 2)


def _next_context(blocks: list[dict], start: int, count: int = 8) -> str:
    parts = []
    for block in blocks[start:start + count]:
        if block["kind"] in {"paragraph", "table"} and len(_one_line(block["text"])) >= 15:
            parts.append(block["text"])
    return _clean("\n".join(parts))


def _split_sentences(text: str) -> list[str]:
    return [
        part.strip()
        for part in re.split(r"(?<=[.!?])\s+(?=[A-Z0-9])", _one_line(text))
        if len(part.strip()) >= 45 and not part.strip().endswith(":")
    ]


def _creative_answer(topic: str, evidence: str) -> str:
    evidence = _trim(evidence, 700)
    return (
        f"{topic.capitalize()} should be positioned around these concrete operating details: "
        f"{evidence} The message should stay specific, partner-facing, and grounded in those "
        "workflows, responsibilities, integrations, or measurement points rather than making broad claims."
    )


def _fact_question(topic: str, heading: str | None = None, key: str | None = None) -> str:
    if key:
        key_text = key.strip().lower()
        if key_text in {"use case", "objective", "goal", "purpose"}:
            return f"What {key_text} is defined for {topic}?"
        return f"What {key_text} is specified for {topic}?"
    if heading:
        clean_heading = re.sub(r"^\d+(?:\.\d+)*\.?\s*", "", heading).strip()
        clean_heading = clean_heading.strip(".: ")
        heading_low = clean_heading.lower()
        if heading_low == "objective":
            return f"What objective is defined for {topic}?"
        if heading_low == "test design":
            return f"How is the test design structured for {topic}?"
        if heading_low == "variants":
            return f"What variants are used for {topic}?"
        if heading_low == "overview":
            return f"What overview is given for {topic}?"
        if "configuration" in heading_low:
            return f"What configuration details control {topic}?"
        if "ad types" in heading_low:
            return f"What ad types are used in {topic}?"
        return f"What role does {clean_heading} play in {topic}?"
    return f"What operational rule is defined for {topic}?"


def _summary_question(topic: str) -> str:
    return f"What are the main objectives, rules, and operating details of {topic}?"


def _reasoning_question(topic: str, text: str) -> str:
    low = text.lower()
    if "enabled" in low and "disabled" in low:
        return f"How does {topic} change when the relevant setting is enabled versus disabled?"
    if "priority" in low or "cpm" in low:
        return f"How does {topic} determine priority when multiple ads or campaigns are eligible?"
    if "workflow" in low or "process" in low or "steps" in low:
        return f"How does the workflow for {topic} move from setup to execution?"
    if "integration" in low:
        return f"How does the integration flow for {topic} connect the required systems?"
    if "measurement" in low or "experiment" in low:
        return f"How does {topic} evaluate performance and decide whether the approach is successful?"
    return f"How do the rules for {topic} work together?"


def _add(candidates: dict, path: Path, question_type: str, question: str, answer: str,
         evidence: str, block: dict):
    question = _one_line(question).rstrip("?. ") + "?"
    answer = _trim(answer)
    if len(answer) < 70 or answer.endswith(":"):
        return
    key = question.lower()
    if key in candidates["seen"] or _norm(answer) in candidates["answers"]:
        return
    candidates["seen"].add(key)
    candidates["answers"].add(_norm(answer))
    record = {
        "id": _make_id(path.name, question, len(candidates["items"]) + 1),
        "question_type": question_type,
        "question": question,
        "answer": answer,
        "ground_truth": answer,
        "source_file": path.name,
        "source_path": str(path),
        "source_type": path.suffix.lower().lstrip("."),
        "evidence": _trim(evidence, 1300),
        "evidence_location": {
            "kind": block.get("kind"),
            "index": block.get("index"),
            "style": block.get("style"),
        },
    }
    if block.get("kind") == "slide":
        record["slide_number"] = block.get("index")
    candidates["items"].append(record)


def _generate_for_file(path: Path, max_per_doc: int) -> tuple[list[dict], dict]:
    blocks, extractor = _extract_docx(path) if path.suffix.lower() == ".docx" else _extract_pptx(path)
    topic = _topic(path)
    candidates = {"items": [], "seen": set(), "answers": set()}
    substantial = [(idx, block) for idx, block in enumerate(blocks) if len(_one_line(block["text"])) >= 60]

    if substantial:
        idx, block = substantial[0]
        answer = _next_context(blocks, idx, 8) if block["kind"] == "paragraph" else block["text"]
        _add(candidates, path, QUESTION_TYPES["summary"], _summary_question(topic), answer, answer, block)
        _add(
            candidates,
            path,
            QUESTION_TYPES["creative"],
            f"Write a concise partner-facing summary that explains {topic}.",
            _creative_answer(topic, answer),
            answer,
            block,
        )

    for idx, block in enumerate(blocks):
        lines = [_one_line(line) for line in block["text"].splitlines() if _one_line(line)]
        if not lines:
            continue
        if block["kind"] == "slide":
            slide_title = next((line for line in lines if _is_heading(line)), lines[0])
            body = "\n".join(line for line in lines if line != slide_title) or block["text"]
            if len(_one_line(body)) < 90:
                continue
            _add(
                candidates,
                path,
                QUESTION_TYPES["summary"],
                f"What should a team understand about {slide_title} in {topic}?",
                body,
                block["text"],
                block,
            )
            continue
        if block["kind"] == "table":
            _add(
                candidates,
                path,
                QUESTION_TYPES["fact"],
                f"What structured information is used for {topic}?",
                block["text"],
                block["text"],
                block,
            )
            continue
        for line_index, line in enumerate(lines):
            if _is_heading(line):
                answer = " ".join(lines[line_index + 1:line_index + 8]) or _next_context(blocks, idx + 1, 8)
                _add(
                    candidates,
                    path,
                    QUESTION_TYPES["fact"],
                    _fact_question(topic, heading=line),
                    answer,
                    block["text"] + "\n" + answer,
                    block,
                )
                break
        for line in lines:
            if ":" in line and 14 <= len(line) <= 280:
                left, right = [part.strip() for part in line.split(":", 1)]
                if 3 <= len(left) <= 70 and len(right) >= 28:
                    if re.match(r"^(section\s+\d+|chapter\s+\d+)\b", left, re.I):
                        continue
                    _add(
                        candidates,
                        path,
                        QUESTION_TYPES["fact"],
                        _fact_question(topic, key=left),
                        line,
                        block["text"],
                        block,
                    )
                    break
        for sentence in _split_sentences(block["text"])[:2]:
            if re.search(
                r"\b(20\d{2}|\d+%|\$|MCC|Octo|Genuin|Mastercard|iHeart|Grubhub|LiveRamp|TikTok|"
                r"campaign|audience|placement|sponsor|feed|ad)\b",
                sentence,
                re.I,
            ):
                _add(
                    candidates,
                    path,
                    QUESTION_TYPES["fact"],
                    f"What specific operational point is stated for {topic}?",
                    sentence,
                    block["text"],
                    block,
                )
                break

    for block in blocks:
        text = _one_line(block["text"])
        if len(text) >= 180 and re.search(
            r"\b(process|workflow|steps|flow|integration|setup|implementation|RACI|responsible|"
            r"approval|priority|prioritization|measurement|experiment)\b",
            text,
            re.I,
        ):
            _add(
                candidates,
                path,
                QUESTION_TYPES["reasoning"],
                _reasoning_question(topic, text),
                text,
                block["text"],
                block,
            )

    grouped = defaultdict(list)
    for item in candidates["items"]:
        grouped[item["question_type"]].append(item)
    selected = []
    order = [
        QUESTION_TYPES["fact"],
        QUESTION_TYPES["summary"],
        QUESTION_TYPES["reasoning"],
        QUESTION_TYPES["creative"],
    ]
    while len(selected) < max_per_doc and any(grouped.values()):
        moved = False
        for question_type in order:
            if grouped[question_type] and len(selected) < max_per_doc:
                selected.append(grouped[question_type].pop(0))
                moved = True
        if not moved:
            break

    context = "\n\n".join(block["text"] for block in blocks)
    stats = {
        "source_file": path.name,
        "source_type": path.suffix.lower().lstrip("."),
        "extractor": extractor,
        "text_blocks_extracted": len(blocks),
        "questions_generated": len(selected),
        "context_chars": len(context),
    }
    return selected, stats


def _write_outputs(subset: str, source_folder: Path, bench_repo: Path | None,
                   output_json: Path, max_per_doc: int):
    files = sorted(
        [
            path for path in source_folder.iterdir()
            if path.is_file()
            and path.suffix.lower() in {".docx", ".pptx"}
            and not path.name.startswith(".~lock")
        ],
        key=lambda path: path.name.lower(),
    )

    questions = []
    corpus_rows = []
    stats = []
    errors = []
    seen_contexts = set()
    for path in files:
        try:
            blocks, _ = _extract_docx(path) if path.suffix.lower() == ".docx" else _extract_pptx(path)
            context = "\n\n".join(block["text"] for block in blocks)
            context_digest = hashlib.sha1(_norm(context).encode("utf-8")).hexdigest()
            if context_digest in seen_contexts:
                continue
            seen_contexts.add(context_digest)
            items, stat = _generate_for_file(path, max_per_doc)
            corpus_name = f"{subset}-{hashlib.sha1(path.name.encode('utf-8')).hexdigest()[:8]}"
            corpus_rows.append({"corpus_name": corpus_name, "context": context})
            for item in items:
                item["source"] = corpus_name
                item["corpus_name"] = corpus_name
            questions.extend(items)
            stats.append(stat | {"corpus_name": corpus_name})
        except Exception as exc:  # noqa: BLE001
            errors.append({"source_file": path.name, "error": f"{type(exc).__name__}: {exc}"})

    for index, item in enumerate(questions, 1):
        item["index"] = index
    deduped = []
    seen_questions = set()
    for item in questions:
        key = item["question"].lower()
        if key in seen_questions:
            continue
        seen_questions.add(key)
        deduped.append(item)
    questions = deduped
    for index, item in enumerate(questions, 1):
        item["index"] = index

    question_rows = [
        {
            "id": item["id"],
            "source": item["source"],
            "question": item["question"],
            "answer": item["ground_truth"],
            "question_type": item["question_type"],
            "evidence": item["evidence"],
            "evidence_triple": item["evidence"],
        }
        for item in questions
    ]
    output_json.parent.mkdir(parents=True, exist_ok=True)
    output_json.write_text(json.dumps(question_rows, indent=2, ensure_ascii=False))

    corpus_path = questions_path = None
    if bench_repo:
        corpus_dir = bench_repo / "Datasets" / "Corpus"
        questions_dir = bench_repo / "Datasets" / "Questions"
        corpus_dir.mkdir(parents=True, exist_ok=True)
        questions_dir.mkdir(parents=True, exist_ok=True)
        corpus_path = corpus_dir / f"{subset}.parquet"
        questions_path = questions_dir / f"{subset}_questions.parquet"
        pd.DataFrame(corpus_rows).to_parquet(corpus_path, index=False)
        pd.DataFrame(question_rows).to_parquet(questions_path, index=False)

    return {
        "json": output_json,
        "corpus_parquet": corpus_path,
        "questions_parquet": questions_path,
        "question_count": len(questions),
        "source_count": len(corpus_rows),
        "errors": errors,
    }


def main(argv=None):
    parser = argparse.ArgumentParser(prog="run.py ref-docs-dataset")
    parser.add_argument(
        "--source-folder",
        default="/home/aditya/Documents/Ref Docs for Knowledge Graph/Ref Docs for Knowledge Graph",
    )
    parser.add_argument("--subset", default="brand_ref")
    parser.add_argument("--max-per-doc", type=int, default=6)
    parser.add_argument(
        "--output-json",
        default=str(Path(config.ROOT) / "data" / "ref_docs_qa" / "ref_docs_questions_answers.json"),
    )
    parser.add_argument(
        "--bench-repo",
        default=None,
        help="GraphRAG-Benchmark repo to receive Datasets/Corpus and Datasets/Questions parquet files",
    )
    parser.add_argument("--no-parquet", action="store_true")
    args = parser.parse_args(argv)

    source_folder = Path(args.source_folder).expanduser().resolve()
    if not source_folder.exists():
        raise SystemExit(f"source folder not found: {source_folder}")
    bench_repo = None
    if not args.no_parquet:
        bench_repo = Path(args.bench_repo or config.BENCH_REPO or "").expanduser().resolve()
        if not bench_repo.exists():
            raise SystemExit("set --bench-repo or BRANDKG_BENCH_REPO, or pass --no-parquet")

    result = _write_outputs(
        subset=args.subset,
        source_folder=source_folder,
        bench_repo=bench_repo,
        output_json=Path(args.output_json).expanduser().resolve(),
        max_per_doc=args.max_per_doc,
    )
    print(f"wrote {result['question_count']} questions from {result['source_count']} files")
    print(f"json: {result['json']}")
    if result["corpus_parquet"]:
        print(f"corpus parquet: {result['corpus_parquet']}")
    if result["questions_parquet"]:
        print(f"questions parquet: {result['questions_parquet']}")
    if result["errors"]:
        print("errors:")
        for error in result["errors"]:
            print(f"  {error['source_file']}: {error['error']}")
    return 0
