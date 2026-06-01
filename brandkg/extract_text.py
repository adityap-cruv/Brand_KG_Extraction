"""Step 1: extract clean text from source documents into data/extracted/*.txt.

Supports .docx, .pptx, .md, .txt. Records a SHA-256 per file in data/manifest.json
so re-runs only re-extract new/changed files (incremental).
"""
from __future__ import annotations
import hashlib
import json
import re
import zipfile
from pathlib import Path

from . import config

SRC = Path(config.SOURCE_DIR)
OUT = config.EXTRACTED_DIR
MANIFEST = config.DATA_DIR / "manifest.json"


def _hash(p: Path) -> str:
    h = hashlib.sha256()
    with p.open("rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return h.hexdigest()


def _docx(p: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(p))
        parts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception:
        # fallback: parse word/document.xml directly (handles corrupt media refs)
        with zipfile.ZipFile(str(p)) as z:
            xml = z.read("word/document.xml").decode("utf-8", "ignore")
        out = []
        for para in re.split(r"</w:p>", xml):
            texts = re.findall(r"<w:t[^>]*>(.*?)</w:t>", para, flags=re.DOTALL)
            line = "".join(texts).strip()
            if line:
                out.append(line)
        return "\n".join(out)


def _pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        sp = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        sp.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        sp.append(" | ".join(cells))
        if len(sp) > 1:
            parts.append("\n".join(sp))
    return "\n\n".join(parts)


def _plain(p: Path) -> str:
    return p.read_text(encoding="utf-8", errors="ignore")


EXTRACTORS = {".docx": _docx, ".pptx": _pptx, ".md": _plain, ".txt": _plain}


def run() -> dict:
    OUT.mkdir(parents=True, exist_ok=True)
    manifest = json.loads(MANIFEST.read_text()) if MANIFEST.exists() else {}
    files = sorted(p for p in SRC.iterdir()
                   if p.is_file() and p.suffix.lower() in EXTRACTORS
                   and not p.name.startswith(("~$", ".~lock")))
    done = skipped = failed = 0
    for p in files:
        h = _hash(p)
        out_path = OUT / f"{p.stem}__{p.suffix.lower().lstrip('.')}.txt"
        rec = manifest.get(p.name)
        if rec and rec.get("hash") == h and out_path.exists():
            skipped += 1
            continue
        try:
            text = EXTRACTORS[p.suffix.lower()](p)
        except Exception as e:
            print(f"  EXTRACT FAIL {p.name}: {e}")
            failed += 1
            continue
        out_path.write_text(text, encoding="utf-8")
        manifest[p.name] = {"hash": h, "out": out_path.name, "chars": len(text)}
        done += 1
        print(f"  OK {p.name} -> {out_path.name} ({len(text):,} chars)")
    MANIFEST.write_text(json.dumps(manifest, indent=2))
    print(f"extract_text: {done} extracted, {skipped} unchanged, {failed} failed")
    return {"extracted": done, "skipped": skipped, "failed": failed}


if __name__ == "__main__":
    run()
